# ============================================================
# pptx_combined.py: 다중 Mermaid 다이어그램을 단일 PPTX로 결합하는 모듈
# 상세: 각 다이어그램을 편집 가능한 네이티브 도형으로 렌더링하여 통합 PPTX 생성
# 생성일: 2026-04-07 | 수정일: 2026-04-07
# ============================================================

from io import BytesIO

from pptx import Presentation
from pptx.util import Inches

from converters.pptx_shapes import mermaid_to_pptx


def create_combined_pptx(diagrams: list[dict]) -> bytes:
    """여러 Mermaid 다이어그램을 슬라이드별로 담은 단일 PPTX 파일을 생성한다.

    각 다이어그램은 편집 가능한 네이티브 도형(사각형, 커넥터, 텍스트)으로
    렌더링되어 PowerPoint에서 바로 편집할 수 있다.

    Args:
        diagrams: ``{"mermaid_code": str, "title": str}`` 딕셔너리의 리스트.

    Returns:
        생성된 PPTX 파일의 바이트 데이터.

    Raises:
        ValueError: diagrams 리스트가 비어 있는 경우.
    """
    if not diagrams:
        raise ValueError("다이어그램 목록이 비어 있습니다.")

    # 첫 번째 다이어그램으로 기본 Presentation 생성
    first_code = diagrams[0].get("mermaid_code", "")
    first_title = diagrams[0].get("title", "Diagram")
    combined_bytes = mermaid_to_pptx(first_code, first_title)
    combined_prs = Presentation(BytesIO(combined_bytes))

    # 나머지 다이어그램의 슬라이드를 하나씩 추가
    for item in diagrams[1:]:
        mermaid_code = item.get("mermaid_code", "")
        title = item.get("title", "Diagram")

        try:
            single_bytes = mermaid_to_pptx(mermaid_code, title)
            single_prs = Presentation(BytesIO(single_bytes))

            # 슬라이드 복사: 원본 슬라이드의 XML을 직접 복제
            for src_slide in single_prs.slides:
                blank_layout = combined_prs.slide_layouts[6]
                new_slide = combined_prs.slides.add_slide(blank_layout)

                # 배경 복사
                new_slide.background.fill.solid()
                from pptx.dml.color import RGBColor
                new_slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                # 모든 도형을 XML 레벨에서 복제
                from lxml import etree
                from pptx.oxml.ns import qn
                src_spTree = src_slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
                dst_spTree = new_slide._element.find(qn("p:cSld")).find(qn("p:spTree"))

                for child in list(src_spTree):
                    tag = child.tag
                    # nvGrpSpPr, grpSpPr는 기본 구조 → 스킵
                    if tag.endswith("}nvGrpSpPr") or tag.endswith("}grpSpPr"):
                        continue
                    dst_spTree.append(etree.fromstring(etree.tostring(child)))

        except Exception:
            pass  # 개별 다이어그램 실패 시 스킵

    output = BytesIO()
    combined_prs.save(output)
    output.seek(0)
    return output.read()
