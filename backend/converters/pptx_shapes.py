# ============================================================
# pptx_shapes.py: Mermaid → 편집 가능한 PowerPoint 도형 변환기
# 상세: Mermaid 코드를 파싱하여 네이티브 도형으로 구성된 PPTX 생성
#       flowchart/graph 및 sequenceDiagram 지원
# 생성일: 2026-04-07 | 수정일: 2026-05-14
# ============================================================

import re
from dataclasses import dataclass, field
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree


# ──────────────────────────────────────────────
# 색상 팔레트 (공통 palette.py에서 가져옴)
# ──────────────────────────────────────────────
from converters.palette import NODE_COLORS, SUBGRAPH_COLORS, TEXT_COLOR, PRIMARY_ACCENT, LINE_COLOR
from converters.text_metrics import estimate_text_size_px


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """'#dbeafe' 형식의 hex 문자열을 RGBColor로 변환."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


_TEXT_RGB = _hex_to_rgb(TEXT_COLOR)
_PRIMARY_ACCENT_RGB = _hex_to_rgb(PRIMARY_ACCENT)
_LINE_RGB = _hex_to_rgb(LINE_COLOR)

_PALETTE = [
    (_hex_to_rgb(fill), _hex_to_rgb(stroke), _TEXT_RGB)
    for fill, stroke in NODE_COLORS
]

_SUBGRAPH_FILLS = [_hex_to_rgb(fill) for fill, _ in SUBGRAPH_COLORS]
_SUBGRAPH_STROKES = [_hex_to_rgb(stroke) for _, stroke in SUBGRAPH_COLORS]
_SUBGRAPH_BORDER = RGBColor(0x94, 0xA3, 0xB8)


# ──────────────────────────────────────────────
# OOXML 보정 헬퍼
# ──────────────────────────────────────────────

def remove_style_element(element):
    """python-pptx가 자동 생성하는 p:style 요소를 제거한다.
    p:style은 테마 색상을 참조하는데, 직접 포매팅과 동시에 존재하면
    PowerPoint가 파일 손상으로 인식한다."""
    style = element.find(qn("p:style"))
    if style is not None:
        element.remove(style)


def _remove_shadow(shape):
    """도형에서 그림자 효과 제거."""
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    for eff in spPr.findall(qn("a:effectLst")):
        spPr.remove(eff)
    etree.SubElement(spPr, qn("a:effectLst"))


def _calc_adj(radius_in, w_in, h_in):
    """원하는 절대 커브 반경(인치)을 adj 값(0~50000)으로 변환."""
    min_dim = min(w_in, h_in)
    if min_dim <= 0:
        return 0
    return int(min(radius_in / min_dim * 50000, 50000))


def _set_corner_radius(shape, adj_val, adj2_val=None):
    """도형의 둥근 모서리 반경을 설정 (OOXML avLst 조정)."""
    prstGeom = shape._element.find('.//' + qn('a:prstGeom'))
    if prstGeom is None:
        return
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is None:
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
    for old in avLst.findall(qn('a:gd')):
        avLst.remove(old)
    prst = prstGeom.get('prst', '')
    if prst == 'roundRect':
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {adj_val}')
    elif prst == 'round2SameRect':
        gd1 = etree.SubElement(avLst, qn('a:gd'))
        gd1.set('name', 'adj1')
        gd1.set('fmla', f'val {adj_val}')
        gd2 = etree.SubElement(avLst, qn('a:gd'))
        gd2.set('name', 'adj2')
        gd2.set('fmla', f'val {adj2_val if adj2_val is not None else 0}')


# ──────────────────────────────────────────────
# 데이터 모델
# ──────────────────────────────────────────────

@dataclass
class Node:
    """Mermaid 노드 정보."""
    id: str
    label: str
    shape: str = "rect"          # rect | round | diamond | circle
    subgraph_id: Optional[str] = None

    # 레이아웃 시 채워짐
    x: float = 0.0               # inches
    y: float = 0.0               # inches
    w: float = 2.0               # inches
    h: float = 0.8               # inches


@dataclass
class Edge:
    """Mermaid 엣지 정보."""
    source: str
    target: str
    label: str = ""
    arrow: str = "-->"           # --> | --- | -.-


@dataclass
class Subgraph:
    """Mermaid 서브그래프(클러스터) 정보."""
    id: str
    label: str
    node_ids: list[str] = field(default_factory=list)

    # 레이아웃 시 채워짐
    x: float = 0.0
    y: float = 0.0
    w: float = 0.0
    h: float = 0.0


@dataclass
class ParsedDiagram:
    """파싱된 Mermaid 다이어그램 전체 구조."""
    nodes: dict[str, Node] = field(default_factory=dict)
    edges: list[Edge] = field(default_factory=list)
    subgraphs: dict[str, Subgraph] = field(default_factory=dict)
    direction: str = "TB"        # TB | LR | RL | BT


# ──────────────────────────────────────────────
# Mermaid 파서
# ──────────────────────────────────────────────

# 노드 정의 패턴 (엣지 선언이 아닌 단독 노드 정의)
_NODE_ALONE_RE = re.compile(
    r"^\s*(?P<id>[A-Za-z0-9_\-]+)"
    r"(?P<shape_open>[\[\(\{\|]+)"
    r"(?P<label>[^\]\)\}\|]*)"
    r"(?P<shape_close>[\]\)\}\|]+)"
    r"\s*$"
)

# 엣지 패턴: A --> B, A -->|label| B, A --- B, A -.-> B 등
_EDGE_RE = re.compile(
    r"^\s*(?P<src>[A-Za-z0-9_\-]+)"
    r"\s*(?P<arrow>--?>|---?\.?-?>?|==+>?|-\.-?>?)"
    r"(?:\|(?P<label>[^\|]*)\|)?"
    r"\s*(?P<dst>[A-Za-z0-9_\-]+)"
    r"\s*$"
)

# 엣지 + 인라인 노드 정의 패턴 (A[label] --> B[label])
_EDGE_WITH_NODES_RE = re.compile(
    r"^\s*(?P<src_id>[A-Za-z0-9_\-]+)"
    r"(?P<src_open>[\[\(\{]+)?"
    r"(?P<src_label>[^\]\)\}]*)?"
    r"(?P<src_close>[\]\)\}]+)?"
    r"\s*(?P<arrow>--?>|---?\.?-?>?|==+>?|-\.-?>?)"
    r"(?:\|(?P<edge_label>[^\|]*)\|)?"
    r"\s*(?P<dst_id>[A-Za-z0-9_\-]+)"
    r"(?P<dst_open>[\[\(\{]+)?"
    r"(?P<dst_label>[^\]\)\}]*)?"
    r"(?P<dst_close>[\]\)\}]+)?"
    r"\s*$"
)


def _shape_from_tokens(open_tok: str, close_tok: str) -> str:
    """괄호 토큰으로 도형 종류를 결정한다."""
    if not open_tok:
        return "rect"
    o = open_tok.strip()
    c = close_tok.strip() if close_tok else ""
    if o == "((":
        return "circle"
    if o == "{":
        return "diamond"
    if o in ("(", "(["):
        return "round"
    return "rect"


def _clean_label(raw: str) -> str:
    """HTML 태그 및 따옴표를 제거하여 순수 텍스트를 반환한다."""
    text = re.sub(r"<[^>]+>", "", raw or "")
    text = text.strip('"').strip("'").strip()
    return text


def parse_mermaid(code: str) -> ParsedDiagram:
    """Mermaid flowchart/graph 코드를 파싱하여 ParsedDiagram을 반환한다."""
    diagram = ParsedDiagram()
    current_subgraph: Optional[Subgraph] = None

    lines = code.splitlines()
    for raw_line in lines:
        line = raw_line.strip()

        # 주석 제거
        line = re.sub(r"%%.*$", "", line).strip()
        if not line:
            continue

        # 방향 선언: graph LR, flowchart TB 등
        dir_match = re.match(
            r"^(?:graph|flowchart)\s+(TB|TD|LR|RL|BT)\s*$", line, re.I
        )
        if dir_match:
            direction = dir_match.group(1).upper()
            diagram.direction = "LR" if direction == "LR" else "TB"
            continue

        # graph/flowchart 선언만 있는 경우 (방향 없음)
        if re.match(r"^(?:graph|flowchart)\s*$", line, re.I):
            continue

        # 서브그래프 시작
        sg_start = re.match(r"^subgraph\s+(?P<id>[A-Za-z0-9_\-]+)\s*(?:\[(?P<label>[^\]]*)\])?\s*$", line, re.I)
        if sg_start:
            sg_id = sg_start.group("id")
            sg_label = _clean_label(sg_start.group("label") or sg_id)
            current_subgraph = Subgraph(id=sg_id, label=sg_label)
            diagram.subgraphs[sg_id] = current_subgraph
            continue

        # 서브그래프 끝
        if re.match(r"^end\s*$", line, re.I):
            current_subgraph = None
            continue

        # 서브그래프 내 direction 무시
        if re.match(r"^direction\s+", line, re.I):
            continue

        # 엣지 + 인라인 노드 파싱 시도
        edge_match = _EDGE_WITH_NODES_RE.match(line)
        if edge_match:
            g = edge_match.groupdict()
            src_id = g["src_id"]
            dst_id = g["dst_id"]

            # 소스 노드 등록
            if src_id not in diagram.nodes:
                shape = _shape_from_tokens(g.get("src_open") or "", g.get("src_close") or "")
                label = _clean_label(g.get("src_label") or src_id)
                diagram.nodes[src_id] = Node(id=src_id, label=label, shape=shape)
            if current_subgraph and src_id not in current_subgraph.node_ids:
                current_subgraph.node_ids.append(src_id)
                diagram.nodes[src_id].subgraph_id = current_subgraph.id

            # 목적지 노드 등록
            if dst_id not in diagram.nodes:
                shape = _shape_from_tokens(g.get("dst_open") or "", g.get("dst_close") or "")
                label = _clean_label(g.get("dst_label") or dst_id)
                diagram.nodes[dst_id] = Node(id=dst_id, label=label, shape=shape)
            if current_subgraph and dst_id not in current_subgraph.node_ids:
                current_subgraph.node_ids.append(dst_id)
                diagram.nodes[dst_id].subgraph_id = current_subgraph.id

            arrow = g.get("arrow") or "-->"
            edge_label = _clean_label(g.get("edge_label") or "")
            diagram.edges.append(Edge(source=src_id, target=dst_id, label=edge_label, arrow=arrow))
            continue

        # 단독 노드 정의
        node_match = _NODE_ALONE_RE.match(line)
        if node_match:
            nid = node_match.group("id")
            shape = _shape_from_tokens(
                node_match.group("shape_open"), node_match.group("shape_close")
            )
            label = _clean_label(node_match.group("label") or nid)
            if nid not in diagram.nodes:
                diagram.nodes[nid] = Node(id=nid, label=label, shape=shape)
            else:
                # 레이블만 업데이트
                diagram.nodes[nid].label = label
                diagram.nodes[nid].shape = shape
            if current_subgraph and nid not in current_subgraph.node_ids:
                current_subgraph.node_ids.append(nid)
                diagram.nodes[nid].subgraph_id = current_subgraph.id

    return diagram


# ──────────────────────────────────────────────
# 레이아웃 엔진
# ──────────────────────────────────────────────

# 슬라이드 치수 (16:9 와이드스크린)
SLIDE_W = 13.333   # inches
SLIDE_H = 7.5      # inches

# 레이아웃 상수
TITLE_H = 0.6      # 제목 영역 높이
MARGIN = 0.3       # 슬라이드 여백
NODE_W = 2.0       # 노드 기본 너비
NODE_H = 0.75      # 노드 기본 높이
H_GAP = 0.45       # 노드 가로 간격
V_GAP = 0.45       # 노드 세로 간격
SG_PAD = 0.3       # 서브그래프 내부 패딩
SG_TITLE_H = 0.35  # 서브그래프 제목 높이


def _layout_nodes_in_grid(
    node_ids: list[str],
    nodes: dict[str, Node],
    start_x: float,
    start_y: float,
    max_width: float,
) -> tuple[float, float]:
    """노드 목록을 그리드로 배치하고, 점유된 (width, height)를 반환한다.

    text_metrics를 사용해 각 노드의 동적 크기를 계산한다 (보강 권고 #2).
    NODE_W / NODE_H 는 fallback 기본값으로만 유지된다.
    """
    if not node_ids:
        return 0.0, 0.0

    # 1) 각 노드의 동적 폭/높이 계산 (CSS px → inches, 96 dpi 기준)
    for nid in node_ids:
        w_px, h_px = estimate_text_size_px(nodes[nid].label, font_size_px=14)
        nodes[nid].w = max(1.4, w_px / 96)
        nodes[nid].h = max(0.55, h_px / 96)

    # 2) 그리드 컬럼 수 결정 (평균 폭 기준)
    avg_w = sum(nodes[nid].w for nid in node_ids) / len(node_ids)
    cols = max(1, int((max_width + H_GAP) / (avg_w + H_GAP)))
    rows = (len(node_ids) + cols - 1) // cols

    # 3) 열별·행별 max 크기 집계 — 보강 권고 #2
    col_widths = [0.0] * cols
    row_heights = [0.0] * rows
    for i, nid in enumerate(node_ids):
        c = i % cols
        r = i // cols
        col_widths[c] = max(col_widths[c], nodes[nid].w)
        row_heights[r] = max(row_heights[r], nodes[nid].h)

    # 4) 좌표 배치
    for i, nid in enumerate(node_ids):
        c = i % cols
        r = i // cols
        nodes[nid].x = start_x + sum(col_widths[:c]) + c * H_GAP
        nodes[nid].y = start_y + sum(row_heights[:r]) + r * V_GAP

    total_w = sum(col_widths) + (cols - 1) * H_GAP
    total_h = sum(row_heights) + (rows - 1) * V_GAP
    return total_w, total_h


def compute_layout(diagram: ParsedDiagram) -> None:
    """파싱된 다이어그램에 좌표를 할당한다 (인플레이스 수정)."""
    avail_w = SLIDE_W - 2 * MARGIN
    avail_h = SLIDE_H - TITLE_H - 2 * MARGIN
    content_x = MARGIN
    content_y = TITLE_H + MARGIN

    # 서브그래프가 있는 경우: 서브그래프 단위로 배치
    if diagram.subgraphs:
        # 서브그래프당 열 수 결정 (최대 3열)
        sg_list = list(diagram.subgraphs.values())
        sg_cols = min(len(sg_list), 3)
        sg_col_w = (avail_w - (sg_cols - 1) * H_GAP) / sg_cols

        cur_x = content_x
        cur_y = content_y
        row_h = 0.0
        col_idx = 0

        for sg in sg_list:
            inner_x = cur_x + SG_PAD
            inner_y = cur_y + SG_PAD + SG_TITLE_H
            inner_w = sg_col_w - 2 * SG_PAD

            node_w, node_h = _layout_nodes_in_grid(
                sg.node_ids, diagram.nodes, inner_x, inner_y, inner_w
            )

            sg.x = cur_x
            sg.y = cur_y
            sg.w = sg_col_w
            sg.h = max(node_h + 2 * SG_PAD + SG_TITLE_H, NODE_H + 2 * SG_PAD + SG_TITLE_H)

            row_h = max(row_h, sg.h)
            col_idx += 1

            if col_idx >= sg_cols:
                cur_x = content_x
                cur_y += row_h + V_GAP
                row_h = 0.0
                col_idx = 0
            else:
                cur_x += sg_col_w + H_GAP

        # 서브그래프에 속하지 않는 노드는 하단에 배치
        orphan_ids = [nid for nid in diagram.nodes if not diagram.nodes[nid].subgraph_id]
        if orphan_ids:
            orphan_y = cur_y + row_h + V_GAP if col_idx > 0 else cur_y
            _layout_nodes_in_grid(orphan_ids, diagram.nodes, content_x, orphan_y, avail_w)

    else:
        # 서브그래프 없음: 전체 노드를 그리드로 배치
        _layout_nodes_in_grid(
            list(diagram.nodes.keys()), diagram.nodes, content_x, content_y, avail_w
        )

    # 오버플로우 감지 → 자동 축소
    _scale_to_fit(diagram, content_x, content_y, avail_w, avail_h)


def _scale_to_fit(
    diagram: "ParsedDiagram",
    content_x: float,
    content_y: float,
    avail_w: float,
    avail_h: float,
) -> None:
    """배치된 다이어그램이 슬라이드를 벗어나면 전체를 축소한다."""
    if not diagram.nodes:
        return

    max_x = max(n.x + n.w for n in diagram.nodes.values())
    max_y = max(n.y + n.h for n in diagram.nodes.values())

    # 서브그래프 영역도 포함
    for sg in diagram.subgraphs.values():
        max_x = max(max_x, sg.x + sg.w)
        max_y = max(max_y, sg.y + sg.h)

    bound_r = content_x + avail_w
    bound_b = content_y + avail_h

    if max_x <= bound_r and max_y <= bound_b:
        return  # 축소 불필요

    scale_x = avail_w / (max_x - content_x) if max_x > bound_r else 1.0
    scale_y = avail_h / (max_y - content_y) if max_y > bound_b else 1.0
    scale = min(scale_x, scale_y)
    scale = max(scale, 0.4)  # 최소 40% — 가독성 보장

    for node in diagram.nodes.values():
        node.x = content_x + (node.x - content_x) * scale
        node.y = content_y + (node.y - content_y) * scale
        node.w *= scale
        node.h *= scale

    for sg in diagram.subgraphs.values():
        sg.x = content_x + (sg.x - content_x) * scale
        sg.y = content_y + (sg.y - content_y) * scale
        sg.w *= scale
        sg.h *= scale


# ──────────────────────────────────────────────
# PPTX 렌더러
# ──────────────────────────────────────────────

def _set_shape_fill(shape, fill_color: RGBColor, stroke_color: RGBColor) -> None:
    """도형의 채우기 색상과 테두리 색상을 설정한다."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color

    line = shape.line
    line.color.rgb = stroke_color
    line.width = Pt(0.75)


def _set_text(shape, text: str, font_size: int = 9, bold: bool = False,
              color: RGBColor = RGBColor(0x1E, 0x29, 0x3B)) -> None:
    """도형의 텍스트 프레임을 설정한다.

    텍스트 길이 대비 도형 폭이 좁으면 폰트를 자동 축소한다 (최소 7pt).
    """
    # 텍스트 길이 기반 폰트 자동 축소
    visual_chars = len(text)
    shape_w_inches = shape.width / 914400  # EMU → inches
    ratio = (shape_w_inches * 96) / max(1, visual_chars * font_size * 0.55)
    effective_size = max(7, int(font_size * min(1.0, ratio)))

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)

    # 기존 단락 초기화
    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER

    run = para.add_run()
    run.text = text
    run.font.size = Pt(effective_size)
    run.font.bold = bold
    run.font.color.rgb = color

    # 한글 폰트 우선: 맑은 고딕
    try:
        rPr = run._r.get_or_add_rPr()
        # 동아시아 폰트 설정
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", "맑은 고딕")
        # 라틴 폰트 설정
        latin = rPr.find(qn("a:latin"))
        if latin is None:
            latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", "맑은 고딕")
    except Exception:
        pass  # 폰트 설정 실패 시 기본값 사용


def _vertical_center_text(shape) -> None:
    """텍스트를 도형 안에서 수직 가운데 정렬한다."""
    from pptx.enum.text import MSO_ANCHOR
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_rounded_rect(slide, x: float, y: float, w: float, h: float,
                      fill: RGBColor, stroke: RGBColor,
                      text: str, font_size: int = 9,
                      text_color: RGBColor = RGBColor(0x1E, 0x29, 0x3B)) -> object:
    """모서리가 둥근 직사각형 도형을 슬라이드에 추가한다."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )

    _set_shape_fill(shape, fill, stroke)
    _set_text(shape, text, font_size=font_size, color=text_color)
    _vertical_center_text(shape)
    remove_style_element(shape._element)
    _remove_shadow(shape)
    return shape


def _add_diamond(slide, x: float, y: float, w: float, h: float,
                 fill: RGBColor, stroke: RGBColor,
                 text: str, font_size: int = 9,
                 text_color: RGBColor = RGBColor(0x1E, 0x29, 0x3B)) -> object:
    """마름모 도형을 슬라이드에 추가한다."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    _set_shape_fill(shape, fill, stroke)
    _set_text(shape, text, font_size=font_size, color=text_color)
    _vertical_center_text(shape)
    remove_style_element(shape._element)
    _remove_shadow(shape)
    return shape


def _add_oval(slide, x: float, y: float, w: float, h: float,
              fill: RGBColor, stroke: RGBColor,
              text: str, font_size: int = 9,
              text_color: RGBColor = RGBColor(0x1E, 0x29, 0x3B)) -> object:
    """타원 도형을 슬라이드에 추가한다."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    _set_shape_fill(shape, fill, stroke)
    _set_text(shape, text, font_size=font_size, color=text_color)
    _vertical_center_text(shape)
    remove_style_element(shape._element)
    _remove_shadow(shape)
    return shape


def _add_connector_elbow(slide, src_shape, dst_shape,
                         label: str = "", dashed: bool = False) -> None:
    """두 도형을 잇는 ELBOW(꺾임선) 커넥터 + 화살표 머리를 추가한다.

    스마트 연결점: 두 도형의 상대 위치에 따라 상/하/좌/우 자동 선택.
    """
    from pptx.enum.shapes import MSO_CONNECTOR

    # 스마트 연결점 선택
    src_cx = src_shape.left + src_shape.width // 2
    src_cy = src_shape.top + src_shape.height // 2
    dst_cx = dst_shape.left + dst_shape.width // 2
    dst_cy = dst_shape.top + dst_shape.height // 2

    dx_abs = abs(dst_cx - src_cx)
    dy_abs = abs(dst_cy - src_cy)

    if dx_abs > dy_abs * 1.5:
        if dst_cx > src_cx:
            sx = src_shape.left + src_shape.width
            sy = src_cy
            ex = dst_shape.left
            ey = dst_cy
        else:
            sx = src_shape.left
            sy = src_cy
            ex = dst_shape.left + dst_shape.width
            ey = dst_cy
    else:
        if dst_cy >= src_cy:
            sx = src_cx
            sy = src_shape.top + src_shape.height
            ex = dst_cx
            ey = dst_shape.top
        else:
            sx = src_cx
            sy = src_shape.top
            ex = dst_cx
            ey = dst_shape.top + dst_shape.height

    # zero extent 방지
    if sx == ex:
        ex += 1
    if sy == ey:
        ey += 1

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW, sx, sy, ex, ey
    )
    connector.line.color.rgb = RGBColor(0x47, 0x55, 0x69)
    connector.line.width = Pt(1.2)

    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # 화살표 머리 추가 (a:tailEnd)
    ln = connector._element.find(".//" + qn("a:ln"))
    if ln is not None:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")

    # OOXML stCxn/endCxn 바인딩 — 도형 이동 시 화살표 추종 활성화 (B 계획)
    def _pick_idx(src, dst):
        """두 도형의 상대 위치로 연결점 인덱스(0=top,1=right,2=bottom,3=left)를 결정한다."""
        s_cx = src.left + src.width // 2
        s_cy = src.top + src.height // 2
        d_cx = dst.left + dst.width // 2
        d_cy = dst.top + dst.height // 2
        if abs(d_cx - s_cx) > abs(d_cy - s_cy):
            return (1, 3) if d_cx > s_cx else (3, 1)
        return (2, 0) if d_cy > s_cy else (0, 2)

    cxnSp = connector._element
    nvCxnSpPr = cxnSp.find(qn("p:nvCxnSpPr"))
    if nvCxnSpPr is None:
        nvCxnSpPr = etree.SubElement(cxnSp, qn("p:nvCxnSpPr"))
    cNvCxnSpPr = nvCxnSpPr.find(qn("p:cNvCxnSpPr"))
    if cNvCxnSpPr is None:
        cNvCxnSpPr = etree.SubElement(nvCxnSpPr, qn("p:cNvCxnSpPr"))

    # 재실행 안전: 기존 stCxn/endCxn 제거
    for tag in ("a:stCxn", "a:endCxn"):
        for old in cNvCxnSpPr.findall(qn(tag)):
            cNvCxnSpPr.remove(old)

    exit_idx, entry_idx = _pick_idx(src_shape, dst_shape)
    stCxn = etree.SubElement(cNvCxnSpPr, qn("a:stCxn"))
    stCxn.set("id", str(src_shape.shape_id))
    stCxn.set("idx", str(exit_idx))
    endCxn = etree.SubElement(cNvCxnSpPr, qn("a:endCxn"))
    endCxn.set("id", str(dst_shape.shape_id))
    endCxn.set("idx", str(entry_idx))

    remove_style_element(connector._element)

    # 엣지 레이블
    if label:
        mx = (sx + ex) // 2
        my = (sy + ey) // 2
        label_box = slide.shapes.add_textbox(
            mx - Inches(0.6), my - Inches(0.15),
            Inches(1.2), Inches(0.3)
        )
        tf = label_box.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = label.replace("\n", " ")
        run.font.size = Pt(7)
        run.font.color.rgb = RGBColor(0x47, 0x55, 0x69)


def _add_node_shape(slide, node: Node, palette_idx: int) -> object:
    """노드의 shape 속성에 맞는 도형을 슬라이드에 추가하고 반환한다."""
    fill_c, stroke_c, text_c = _PALETTE[palette_idx % len(_PALETTE)]

    if node.shape == "diamond":
        return _add_diamond(
            slide, node.x, node.y, node.w, node.h,
            fill_c, stroke_c, node.label, text_color=text_c
        )
    elif node.shape == "circle":
        return _add_oval(
            slide, node.x, node.y, node.w, node.h,
            fill_c, stroke_c, node.label, text_color=text_c
        )
    else:
        # rect, round 모두 rounded rectangle로 렌더링
        return _add_rounded_rect(
            slide, node.x, node.y, node.w, node.h,
            fill_c, stroke_c, node.label, text_color=text_c
        )


def _node_center(node: Node) -> tuple[float, float]:
    """노드의 중심 좌표를 반환한다."""
    return node.x + node.w / 2, node.y + node.h / 2


def _add_subgraph_box(slide, sg: Subgraph, fill_color: RGBColor, idx: int) -> None:
    """서브그래프를 draw.io 스타일 2개 도형으로 추가한다.
    1) 컨테이너: ROUNDED_RECTANGLE, 흰색 배경 + 얇은 윤곽선, 그림자 없음
    2) 제목: ROUND_2_SAME_RECTANGLE, 연한 색상 배경 + 텍스트 직접 포함
    """
    stroke_color = _SUBGRAPH_STROKES[idx % len(_SUBGRAPH_STROKES)]
    title_h = 0.26
    CORNER_RADIUS_IN = 0.10
    line_w = Pt(0.75)

    adj_container = _calc_adj(CORNER_RADIUS_IN, sg.w, sg.h)
    adj_title = _calc_adj(CORNER_RADIUS_IN, sg.w, title_h)

    # 1) 컨테이너 — 흰색 배경 + 얇은 윤곽선
    container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(sg.x), Inches(sg.y), Inches(sg.w), Inches(sg.h)
    )
    container.fill.solid()
    container.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    container.line.color.rgb = stroke_color
    container.line.width = line_w
    _set_corner_radius(container, adj_container)
    _remove_shadow(container)
    remove_style_element(container._element)
    tf = container.text_frame
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    # 2) 제목 도형 — 둥근 위쪽 모서리 + 텍스트 직접 포함
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
        Inches(sg.x), Inches(sg.y), Inches(sg.w), Inches(title_h)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = fill_color
    title_shape.line.color.rgb = stroke_color
    title_shape.line.width = line_w
    _set_corner_radius(title_shape, adj_title, adj2_val=0)
    _remove_shadow(title_shape)
    remove_style_element(title_shape._element)

    # 제목 텍스트를 도형 안에 직접 배치
    tf2 = title_shape.text_frame
    tf2.auto_size = None
    tf2.word_wrap = True
    tf2.margin_top = Pt(3)
    tf2.margin_bottom = Pt(2)
    tf2.margin_left = Pt(6)
    tf2.margin_right = Pt(6)
    txBody = title_shape._element.find(qn("p:txBody"))
    if txBody is not None:
        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sg.label
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = stroke_color

    try:
        rPr = run._r.get_or_add_rPr()
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", "맑은 고딕")
        latin = rPr.find(qn("a:latin"))
        if latin is None:
            latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", "맑은 고딕")
    except Exception:
        pass


# ──────────────────────────────────────────────
# 시퀀스 다이어그램 파서 / 렌더러
# ──────────────────────────────────────────────

# 시퀀스 다이어그램 레이아웃 상수
_SEQ_PARTICIPANT_W = 2.0    # 참여자 박스 너비 (inches)
_SEQ_PARTICIPANT_H = 0.6    # 참여자 박스 높이 (inches)
_SEQ_PARTICIPANT_GAP = 0.5  # 참여자 간격 (inches)
_SEQ_MSG_GAP = 0.5          # 메시지 세로 간격 (inches)
_SEQ_TOP_MARGIN = 1.0       # 상단 여백 (제목 포함, inches)

# 시퀀스 메시지 화살표 패턴
_SEQ_ARROW_RE = re.compile(
    r"^\s*(?P<src>[A-Za-z0-9_]+)\s*"
    r"(?P<arrow>-->>|--?>|-->|->|->>)"
    r"\s*(?P<dst>[A-Za-z0-9_]+)\s*"
    r":\s*(?P<label>.+)$"
)

# participant/actor 선언 패턴
_SEQ_PARTICIPANT_RE = re.compile(
    r"^\s*(?:participant|actor)\s+(?P<id>[A-Za-z0-9_]+)"
    r"(?:\s+as\s+(?P<label>.+))?\s*$",
    re.I
)


def _parse_sequence(mermaid_code: str) -> tuple[list[tuple[str, str]], list[tuple[str, str, str, bool]]]:
    """시퀀스 다이어그램 Mermaid 코드를 파싱한다.

    Returns:
        (participants, messages) 튜플.
        participants: [(id, label), ...] 순서 보장 리스트.
        messages: [(source, target, label, dashed), ...] 리스트.
    """
    participants: list[tuple[str, str]] = []
    participant_ids: list[str] = []
    messages: list[tuple[str, str, str, bool]] = []

    for raw_line in mermaid_code.splitlines():
        line = raw_line.strip()
        if not line or line.lower().startswith("sequencediagram"):
            continue

        # participant/actor 선언
        p_match = _SEQ_PARTICIPANT_RE.match(line)
        if p_match:
            pid = p_match.group("id")
            plabel = re.sub(r"<br\s*/?>", " ", (p_match.group("label") or pid)).strip()
            if pid not in participant_ids:
                participant_ids.append(pid)
                participants.append((pid, plabel))
            continue

        # 메시지 화살표
        m_match = _SEQ_ARROW_RE.match(line)
        if m_match:
            src = m_match.group("src")
            dst = m_match.group("dst")
            label = re.sub(r"<br\s*/?>", " ", m_match.group("label")).strip()
            arrow = m_match.group("arrow")
            dashed = "--" in arrow
            messages.append((src, dst, label, dashed))

            # 암묵적 참여자 등록
            for actor_id in (src, dst):
                if actor_id not in participant_ids:
                    participant_ids.append(actor_id)
                    participants.append((actor_id, actor_id))

    return participants, messages


def _render_sequence(mermaid_code: str, title: str = "") -> bytes:
    """시퀀스 다이어그램 Mermaid 코드를 네이티브 도형 PPTX로 변환한다.

    제어 프레임(alt/opt/loop/par/critical/break), Note 박스, 라벨 클램프 포함.

    Args:
        mermaid_code: sequenceDiagram Mermaid 코드 문자열.
        title: 슬라이드 상단 제목.

    Returns:
        생성된 PPTX 파일의 바이트 데이터.
    """
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from converters.drawio import _parse_sequence_events

    participants, _messages_unused = _parse_sequence(mermaid_code)
    events = _parse_sequence_events(mermaid_code)

    if not participants:
        raise ValueError("시퀀스 다이어그램에서 참여자를 찾을 수 없습니다.")

    # PPTX 생성
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 배경: 흰색
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 제목
    if title:
        title_box = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(0.1),
            Inches(SLIDE_W - 2 * MARGIN), Inches(TITLE_H)
        )
        tf = title_box.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = title
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        try:
            rPr = run._r.get_or_add_rPr()
            ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", "맑은 고딕")
            latin = rPr.find(qn("a:latin"))
            if latin is None:
                latin = etree.SubElement(rPr, qn("a:latin"))
            latin.set("typeface", "맑은 고딕")
        except Exception:
            pass

        # 제목 하단 구분선
        line_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(MARGIN), Inches(TITLE_H + 0.05),
            Inches(SLIDE_W - 2 * MARGIN), Inches(0.03)
        )
        line_bar.fill.solid()
        line_bar.fill.fore_color.rgb = _PRIMARY_ACCENT_RGB
        line_bar.line.fill.background()
        remove_style_element(line_bar._element)
        _remove_shadow(line_bar)

    # ── 스케일 계산 ──────────────────────────────
    n_part = len(participants)
    part_y = _SEQ_TOP_MARGIN

    avail_w = SLIDE_W - 2 * MARGIN
    nat_w = n_part * _SEQ_PARTICIPANT_W + max(n_part - 1, 0) * _SEQ_PARTICIPANT_GAP
    sx_scale = min(1.0, avail_w / nat_w) if nat_w else 1.0
    pw = _SEQ_PARTICIPANT_W * sx_scale
    pgap = _SEQ_PARTICIPANT_GAP * sx_scale
    ph = _SEQ_PARTICIPANT_H
    pfont = max(6, int(round(9 * (sx_scale ** 0.5))))

    msg_y_start = part_y + ph + 0.3
    avail_h = SLIDE_H - MARGIN - msg_y_start - ph  # 하단 참여자 박스 공간 확보

    # 논리 라인 수: msg + note 각각 한 줄씩 (group_start/end는 y 소비 없음)
    n_logical = sum(1 for e in events if e["kind"] in ("msg", "note"))
    nat_msg_h = max(n_logical, 1) * _SEQ_MSG_GAP + 0.2
    sy_scale = min(1.0, avail_h / nat_msg_h) if nat_msg_h else 1.0
    mgap = _SEQ_MSG_GAP * sy_scale

    total_part_w = n_part * pw + max(n_part - 1, 0) * pgap
    start_x = max(MARGIN, (SLIDE_W - total_part_w) / 2.0)

    # 참여자 중심 x 맵
    part_centers: dict[str, float] = {}
    for i, (pid, _) in enumerate(participants):
        px = start_x + i * (pw + pgap)
        part_centers[pid] = px + pw / 2.0

    # ── Phase 1: Y 좌표 사전 계산 + 프레임 정보 수집 ──
    y_cursor = msg_y_start
    group_stack: list[dict] = []
    frame_queue: list[dict] = []
    event_data: list[dict] = []   # {ev, y} for msg/note events

    for ev in events:
        kind = ev["kind"]
        if kind == "msg":
            event_data.append({"ev": ev, "y": y_cursor})
            y_cursor += mgap
        elif kind == "note":
            event_data.append({"ev": ev, "y": y_cursor})
            y_cursor += mgap
        elif kind == "group_start":
            g = {"type": ev["type"], "label": ev["label"],
                 "y_top": y_cursor, "else_ys": []}
            group_stack.append(g)
            event_data.append({"ev": ev, "y": y_cursor})
        elif kind == "group_else":
            if group_stack:
                group_stack[-1]["else_ys"].append(y_cursor)
            event_data.append({"ev": ev, "y": y_cursor})
        elif kind == "group_end":
            if group_stack:
                g = group_stack.pop()
                g["y_bot"] = y_cursor
                frame_queue.append(g)
            event_data.append({"ev": ev, "y": y_cursor})

    msg_y_end = y_cursor + 0.2 * sy_scale

    # 프레임 x 경계 (모든 참여자 포함)
    _FRAME_FILL = RGBColor(0xEF, 0xF6, 0xFF)
    _FRAME_STROKE = RGBColor(0x60, 0xA5, 0xFA)
    _NOTE_FILL = RGBColor(0xFF, 0xFB, 0xD5)
    _NOTE_STROKE = RGBColor(0xD9, 0x7F, 0x0E)
    _FRAME_PAD_Y = 0.14

    frame_x = max(MARGIN + 0.02, start_x - 0.05)
    frame_right = min(
        SLIDE_W - MARGIN - 0.02,
        start_x + (n_part - 1) * (pw + pgap) + pw + 0.05
    )
    frame_w = max(0.3, frame_right - frame_x)

    # ── Phase 2: 프레임 먼저 그리기 (화살표 뒤 레이어) ──
    for frame in frame_queue:
        gtype = frame["type"]
        glabel = frame.get("label", "")
        gy_top = frame["y_top"]
        gy_bot = frame.get("y_bot", gy_top + mgap)
        else_ys = frame.get("else_ys", [])

        fy = gy_top - _FRAME_PAD_Y
        fh = max(0.3, gy_bot - gy_top + _FRAME_PAD_Y + 0.08)

        frame_sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(frame_x), Inches(fy),
            Inches(frame_w), Inches(fh)
        )
        frame_sh.fill.solid()
        frame_sh.fill.fore_color.rgb = _FRAME_FILL
        frame_sh.line.color.rgb = _FRAME_STROKE
        frame_sh.line.width = Pt(1.0)
        frame_sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        frame_sh.text_frame.text = ""
        remove_style_element(frame_sh._element)
        _remove_shadow(frame_sh)

        # 좌상단 라벨 탭
        tab_text = gtype
        if glabel:
            tab_text += f" {glabel}"
        tab_w = max(0.5, min(len(tab_text) * 0.075 + 0.1, frame_w * 0.55, 2.5))
        tab_h = 0.2
        tab_sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(frame_x), Inches(fy),
            Inches(tab_w), Inches(tab_h)
        )
        tab_sh.fill.solid()
        tab_sh.fill.fore_color.rgb = _FRAME_STROKE
        tab_sh.line.fill.background()
        remove_style_element(tab_sh._element)
        _remove_shadow(tab_sh)
        tf_tab = tab_sh.text_frame
        tf_tab.margin_top = Pt(1)
        tf_tab.margin_bottom = Pt(1)
        tf_tab.margin_left = Pt(3)
        tf_tab.margin_right = Pt(3)
        tf_tab.clear()
        p_tab = tf_tab.paragraphs[0]
        p_tab.alignment = PP_ALIGN.LEFT
        r_tab = p_tab.add_run()
        r_tab.text = tab_text[:40]
        r_tab.font.size = Pt(6)
        r_tab.font.bold = True
        r_tab.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # else 구분선
        for ey in else_ys:
            l_emu = Inches(frame_x)
            r_emu = Inches(frame_x + frame_w)
            y_emu = Inches(ey)
            if l_emu == r_emu:
                r_emu += 1
            div = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                l_emu, y_emu, r_emu, y_emu
            )
            div.line.color.rgb = _FRAME_STROKE
            div.line.width = Pt(0.75)
            div.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            remove_style_element(div._element)
            _remove_shadow(div)

    # ── Phase 3: 참여자 박스 (상단 + 하단 반복, mermaid 스타일) ──
    def _draw_participant_row(box_y: float) -> None:
        for i, (pid, plabel) in enumerate(participants):
            px = start_x + i * (pw + pgap)
            fill_c, stroke_c, text_c = _PALETTE[i % len(_PALETTE)]
            shape = _add_rounded_rect(
                slide, px, box_y, pw, ph,
                fill_c, stroke_c, plabel, font_size=pfont, text_color=text_c
            )
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.bold = True

    _draw_participant_row(part_y)

    # ── Phase 4: 생명선 (수직 점선): 상단 박스 하단 → 하단 박스 상단 ──
    for pid, cx in part_centers.items():
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(cx), Inches(part_y + ph),
            Inches(cx) + 1, Inches(msg_y_end)
        )
        connector.line.color.rgb = _LINE_RGB
        connector.line.width = Pt(0.75)
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        remove_style_element(connector._element)
        _remove_shadow(connector)

    # 하단 참여자 박스 반복 (생명선 끝)
    _draw_participant_row(msg_y_end)

    # ── Phase 5: 이벤트 렌더 (메시지 화살표 + Note 박스) ──
    lbl_w = min(1.6, pw * 1.4)     # 메시지 라벨 텍스트박스 폭
    loop_offset_raw = 0.3 * sx_scale  # self-loop 가로 돌출량

    for ed in event_data:
        ev = ed["ev"]
        ey_pos = ed["y"]
        kind = ev["kind"]

        if kind == "msg":
            src = ev["src"]
            dst = ev["dst"]
            label = ev["label"]
            dashed = ev["dashed"]

            if src not in part_centers or dst not in part_centers:
                continue

            msg_y_emu = Inches(ey_pos)
            sx_emu = Inches(part_centers[src])
            dx_emu = Inches(part_centers[dst])

            if src == dst:
                # 자기 자신 메시지: 클램프된 오른쪽 루프
                cx_in = part_centers[src]
                right_margin = SLIDE_W - MARGIN - cx_in
                loop_offset = max(0.1, min(loop_offset_raw, right_margin - 0.05))
                loop_h_in = mgap * 0.4
                lo_emu = Inches(loop_offset)
                lh_emu = Inches(loop_h_in)

                h_conn = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    sx_emu, msg_y_emu,
                    sx_emu + lo_emu, msg_y_emu
                )
                h_conn.line.color.rgb = RGBColor(0x47, 0x55, 0x69)
                h_conn.line.width = Pt(1.2)
                if dashed:
                    h_conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                remove_style_element(h_conn._element)
                _remove_shadow(h_conn)

                v_conn = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    sx_emu + lo_emu, msg_y_emu,
                    sx_emu + lo_emu, msg_y_emu + lh_emu
                )
                v_conn.line.color.rgb = RGBColor(0x47, 0x55, 0x69)
                v_conn.line.width = Pt(1.2)
                if dashed:
                    v_conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                remove_style_element(v_conn._element)
                _remove_shadow(v_conn)

                r_conn = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    sx_emu + lo_emu, msg_y_emu + lh_emu,
                    sx_emu, msg_y_emu + lh_emu
                )
                r_conn.line.color.rgb = RGBColor(0x47, 0x55, 0x69)
                r_conn.line.width = Pt(1.2)
                if dashed:
                    r_conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                ln = r_conn._element.find(".//" + qn("a:ln"))
                if ln is not None:
                    tail = etree.SubElement(ln, qn("a:tailEnd"))
                    tail.set("type", "triangle")
                    tail.set("w", "med")
                    tail.set("len", "med")
                remove_style_element(r_conn._element)
                _remove_shadow(r_conn)

                if label:
                    lx = cx_in + loop_offset
                    lx = max(MARGIN, min(SLIDE_W - MARGIN - lbl_w, lx))
                    txb = slide.shapes.add_textbox(
                        Inches(lx),
                        msg_y_emu - Inches(0.12),
                        Inches(lbl_w), Inches(0.25)
                    )
                    tf = txb.text_frame
                    tf.word_wrap = False
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = label
                    run.font.size = Pt(7)
                    run.font.color.rgb = RGBColor(0x47, 0x55, 0x69)

            else:
                # 일반 수평 메시지 화살표
                if sx_emu == dx_emu:
                    dx_emu += 1

                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    sx_emu, msg_y_emu,
                    dx_emu, msg_y_emu
                )
                connector.line.color.rgb = RGBColor(0x47, 0x55, 0x69)
                connector.line.width = Pt(1.2)
                if dashed:
                    connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                ln = connector._element.find(".//" + qn("a:ln"))
                if ln is not None:
                    tail = etree.SubElement(ln, qn("a:tailEnd"))
                    tail.set("type", "triangle")
                    tail.set("w", "med")
                    tail.set("len", "med")
                remove_style_element(connector._element)
                _remove_shadow(connector)

                if label:
                    sx_in = part_centers[src]
                    dx_in = part_centers[dst]
                    mx_in = (sx_in + dx_in) / 2
                    lx = mx_in - lbl_w / 2
                    lx = max(MARGIN, min(SLIDE_W - MARGIN - lbl_w, lx))
                    txb = slide.shapes.add_textbox(
                        Inches(lx),
                        msg_y_emu - Inches(0.18),
                        Inches(lbl_w), Inches(0.3)
                    )
                    tf = txb.text_frame
                    tf.word_wrap = False
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = label
                    run.font.size = Pt(7)
                    run.font.color.rgb = RGBColor(0x47, 0x55, 0x69)

        elif kind == "note":
            actors = ev["actors"]
            text = ev["text"]

            actor_xs = [part_centers[a] for a in actors if a in part_centers]
            if not actor_xs:
                actor_xs = list(part_centers.values())[:1] if part_centers else [start_x + pw / 2]

            note_left_x = min(actor_xs) - pw / 2
            note_right_x = max(actor_xs) + pw / 2
            note_w = max(1.0, min(note_right_x - note_left_x, SLIDE_W - 2 * MARGIN))
            note_cx = (min(actor_xs) + max(actor_xs)) / 2
            note_x = note_cx - note_w / 2
            note_x = max(MARGIN, min(SLIDE_W - MARGIN - note_w, note_x))

            note_h = 0.28
            note_y = ey_pos - note_h / 2

            note_sh = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(note_x), Inches(note_y),
                Inches(note_w), Inches(note_h)
            )
            note_sh.fill.solid()
            note_sh.fill.fore_color.rgb = _NOTE_FILL
            note_sh.line.color.rgb = _NOTE_STROKE
            note_sh.line.width = Pt(0.75)
            remove_style_element(note_sh._element)
            _remove_shadow(note_sh)

            tf_note = note_sh.text_frame
            tf_note.word_wrap = True
            tf_note.margin_top = Pt(2)
            tf_note.margin_bottom = Pt(2)
            tf_note.margin_left = Pt(4)
            tf_note.margin_right = Pt(4)
            tf_note.clear()
            p_note = tf_note.paragraphs[0]
            p_note.alignment = PP_ALIGN.CENTER
            r_note = p_note.add_run()
            r_note.text = text
            r_note.font.size = Pt(7)
            r_note.font.color.rgb = RGBColor(0x78, 0x35, 0x00)

    # BytesIO로 저장 후 바이트 반환
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


# ──────────────────────────────────────────────
# erDiagram 테이블 렌더러
# ──────────────────────────────────────────────

def _render_er_diagram(mermaid_code: str, title: str = "") -> bytes:
    """erDiagram Mermaid 코드를 엔티티 테이블 + 관계 커넥터 PPTX로 변환한다.

    각 엔티티를 헤더(이름) + 본문(속성 목록) 형태의 표로 렌더링하고,
    관계를 ELBOW 커넥터 + 카디널리티 라벨로 연결한다.
    """
    from pptx.enum.text import MSO_ANCHOR
    from converters.drawio import _parse_er, _parse_er_attrs

    entities, relations = _parse_er(mermaid_code)
    attrs_map = _parse_er_attrs(mermaid_code)

    if not entities:
        raise ValueError("erDiagram에서 엔티티를 찾을 수 없습니다.")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    content_y = MARGIN
    if title:
        title_box = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(0.1),
            Inches(SLIDE_W - 2 * MARGIN), Inches(TITLE_H)
        )
        tf = title_box.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = title
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        try:
            rPr = run._r.get_or_add_rPr()
            ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", "맑은 고딕")
            latin = rPr.find(qn("a:latin"))
            if latin is None:
                latin = etree.SubElement(rPr, qn("a:latin"))
            latin.set("typeface", "맑은 고딕")
        except Exception:
            pass
        line_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(MARGIN), Inches(TITLE_H + 0.05),
            Inches(SLIDE_W - 2 * MARGIN), Inches(0.03)
        )
        line_bar.fill.solid()
        line_bar.fill.fore_color.rgb = _PRIMARY_ACCENT_RGB
        line_bar.line.fill.background()
        remove_style_element(line_bar._element)
        _remove_shadow(line_bar)
        content_y = TITLE_H + MARGIN

    avail_w = SLIDE_W - 2 * MARGIN
    avail_h = SLIDE_H - content_y - MARGIN

    # 테이블 치수 상수
    HEADER_H = 0.38
    ROW_H = 0.24
    TBL_W_BASE = 2.0
    TBL_W_GAP = 0.35
    TBL_H_GAP = 0.30

    n = len(entities)
    # 열 수: 슬라이드 너비에 맞춰 자동 결정
    cols = max(1, min(n, int((avail_w + TBL_W_GAP) / (TBL_W_BASE + TBL_W_GAP))))
    rows_count = (n + cols - 1) // cols

    # 각 엔티티 높이 계산
    entity_heights = {
        ent: HEADER_H + len(attrs_map.get(ent, [])) * ROW_H
        for ent in entities
    }

    # 그리드 행별 최대 높이
    grid_row_h = []
    for r in range(rows_count):
        row_ents = entities[r * cols:(r + 1) * cols]
        grid_row_h.append(max((entity_heights[e] for e in row_ents), default=HEADER_H))

    total_w = cols * TBL_W_BASE + (cols - 1) * TBL_W_GAP
    total_h = sum(grid_row_h) + (rows_count - 1) * TBL_H_GAP

    # scale-to-fit
    scale = 1.0
    if total_w > avail_w:
        scale = min(scale, avail_w / total_w)
    if total_h > avail_h:
        scale = min(scale, avail_h / total_h)
    scale = max(scale, 0.35)

    tbl_w = TBL_W_BASE * scale
    tbl_w_gap = TBL_W_GAP * scale
    tbl_h_gap = TBL_H_GAP * scale
    header_h = HEADER_H * scale
    row_h_s = ROW_H * scale

    # 중앙 정렬
    actual_w = cols * tbl_w + (cols - 1) * tbl_w_gap
    actual_h = sum(h * scale for h in grid_row_h) + (rows_count - 1) * tbl_h_gap
    start_x = MARGIN + (avail_w - actual_w) / 2
    start_y = max(content_y, content_y + (avail_h - actual_h) / 2)

    font_hdr = max(6, int(9 * scale))
    font_body = max(5, int(7 * scale))

    # 엔티티 rect 정보: {ent: (x, y, w, h)} — 커넥터용
    entity_rects: dict[str, tuple] = {}

    from pptx.enum.text import MSO_ANCHOR as _MSO_ANCHOR

    def _style_table_grid(tbl) -> None:
        """네이티브 표를 'No Style, Table Grid'(얇은 격자선, 밴딩 없음)로 설정."""
        tblPr = tbl._tbl.tblPr
        if tblPr is None:
            return
        tblPr.set("firstRow", "0")
        tblPr.set("bandRow", "0")
        sid = tblPr.find(qn("a:tableStyleId"))
        if sid is None:
            sid = etree.SubElement(tblPr, qn("a:tableStyleId"))
        sid.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"

    def _set_cell(cell, text, size, bold, color, fill, align) -> None:
        """표 셀 텍스트/색상/정렬 설정."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)
        cell.margin_left = Pt(4)
        cell.margin_right = Pt(3)
        cell.vertical_anchor = _MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = False
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    for i, ent in enumerate(entities):
        col = i % cols
        row_idx = i // cols

        tx = start_x + col * (tbl_w + tbl_w_gap)
        ty = (start_y
              + sum(grid_row_h[r] * scale for r in range(row_idx))
              + row_idx * tbl_h_gap)

        attrs = attrs_map.get(ent, [])
        ent_h = header_h + len(attrs) * row_h_s

        fill_c, stroke_c, text_c = _PALETTE[i % len(_PALETTE)]

        # ── 네이티브 표로 DB 테이블 형상화 (헤더 행 + 컬럼별 type|name 행) ──
        n_rows = 1 + len(attrs)
        gf = slide.shapes.add_table(
            n_rows, 2,
            Inches(tx), Inches(ty), Inches(tbl_w), Inches(ent_h)
        )
        tbl = gf.table
        _style_table_grid(tbl)

        # 열 폭 (type 42% / name 58%) · 행 높이
        tbl.columns[0].width = Inches(tbl_w * 0.42)
        tbl.columns[1].width = Inches(tbl_w * 0.58)
        tbl.rows[0].height = Inches(header_h)
        for rr in range(1, n_rows):
            tbl.rows[rr].height = Inches(row_h_s)

        # 헤더 행: 2열 병합 + 엔티티명 (팔레트 색 채움)
        hc = tbl.cell(0, 0)
        hc.merge(tbl.cell(0, 1))
        _set_cell(hc, ent, font_hdr, True, text_c, fill_c, PP_ALIGN.CENTER)

        # 본문 행: type | name (흰 배경)
        _WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        _TYPE_C = RGBColor(0x47, 0x55, 0x69)
        _NAME_C = RGBColor(0x1E, 0x29, 0x3B)
        for k, (atype, aname) in enumerate(attrs):
            _set_cell(tbl.cell(k + 1, 0), atype, font_body, False, _TYPE_C, _WHITE, PP_ALIGN.LEFT)
            _set_cell(tbl.cell(k + 1, 1), aname, font_body, False, _NAME_C, _WHITE, PP_ALIGN.LEFT)

        entity_rects[ent] = (tx, ty, tbl_w, ent_h)

    # 관계 커넥터 + 카디널리티 라벨
    for (a, b, lbl) in relations:
        if a not in entity_rects or b not in entity_rects:
            continue
        ax, ay, aw, ah = entity_rects[a]
        bx, by, bw, bh = entity_rects[b]

        acx, acy = ax + aw / 2, ay + ah / 2
        bcx, bcy = bx + bw / 2, by + bh / 2
        dx = bcx - acx
        dy = bcy - acy

        if abs(dx) >= abs(dy):
            if dx >= 0:
                sx, sy = ax + aw, acy
                ex, ey = bx, bcy
            else:
                sx, sy = ax, acy
                ex, ey = bx + bw, bcy
        else:
            if dy >= 0:
                sx, sy = acx, ay + ah
                ex, ey = bcx, by
            else:
                sx, sy = acx, ay
                ex, ey = bcx, by + bh

        sx_emu = Inches(sx)
        sy_emu = Inches(sy)
        ex_emu = Inches(ex)
        ey_emu = Inches(ey)
        if sx_emu == ex_emu:
            ex_emu += 1
        if sy_emu == ey_emu:
            ey_emu += 1

        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW, sx_emu, sy_emu, ex_emu, ey_emu
        )
        conn.line.color.rgb = RGBColor(0x47, 0x55, 0x69)
        conn.line.width = Pt(1.0)
        remove_style_element(conn._element)
        _remove_shadow(conn)

        if lbl:
            mx = (sx + ex) / 2
            my = (sy + ey) / 2
            lw = min(1.6, max(0.5, len(lbl) * 0.07))
            lbl_x = max(MARGIN, min(SLIDE_W - MARGIN - lw, mx - lw / 2))
            lbl_box = slide.shapes.add_textbox(
                Inches(lbl_x), Inches(my - 0.12),
                Inches(lw), Inches(0.25)
            )
            tf_l = lbl_box.text_frame
            tf_l.word_wrap = False
            p_l = tf_l.paragraphs[0]
            p_l.alignment = PP_ALIGN.CENTER
            r_l = p_l.add_run()
            r_l.text = lbl.strip()
            r_l.font.size = Pt(6)
            r_l.font.color.rgb = RGBColor(0x47, 0x55, 0x69)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


# ──────────────────────────────────────────────
# 공개 API
# ──────────────────────────────────────────────

def mermaid_to_pptx(mermaid_code: str, title: str = "") -> bytes:
    """Mermaid flowchart/graph 코드를 네이티브 도형 PPTX로 변환한다.

    모든 도형은 PowerPoint에서 개별 선택·이동·편집이 가능하다.

    Args:
        mermaid_code: 변환할 Mermaid 다이어그램 코드 문자열.
        title: 슬라이드 상단에 표시할 제목. 빈 문자열이면 제목 영역을 생략한다.

    Returns:
        생성된 PPTX 파일의 바이트 데이터 (BytesIO 기반).

    Raises:
        ValueError: Mermaid 코드가 비어 있거나 노드를 찾을 수 없는 경우.
    """
    if not mermaid_code or not mermaid_code.strip():
        raise ValueError("Mermaid 코드가 비어 있습니다.")

    # 타입 판별: %% 주석·빈 줄을 건너뛴 첫 의미 있는 줄 기준
    from converters.drawio import _first_meaningful_line
    first_line = _first_meaningful_line(mermaid_code)

    # 시퀀스 다이어그램 분기
    if 'sequencediagram' in first_line.lower().replace(' ', ''):
        return _render_sequence(mermaid_code, title)

    # erDiagram 분기 — 엔티티를 컬럼 테이블로, 관계를 커넥터로 렌더링
    if first_line.startswith('erDiagram'):
        return _render_er_diagram(mermaid_code, title)

    # 1. 파싱 (flowchart/graph)
    diagram = parse_mermaid(mermaid_code)

    if not diagram.nodes:
        raise ValueError("Mermaid 코드에서 노드를 찾을 수 없습니다.")

    # 2. 레이아웃 계산
    compute_layout(diagram)

    # 3. PPTX 생성
    prs = Presentation()

    # 16:9 슬라이드 크기 설정
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 빈 레이아웃 사용 (인덱스 6 = blank layout)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 4. 슬라이드 배경: 흰색
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 5. 제목 추가
    if title:
        title_box = slide.shapes.add_textbox(
            Inches(MARGIN),
            Inches(0.1),
            Inches(SLIDE_W - 2 * MARGIN),
            Inches(TITLE_H)
        )
        tf = title_box.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = title
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

        try:
            rPr = run._r.get_or_add_rPr()
            ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", "맑은 고딕")
            latin = rPr.find(qn("a:latin"))
            if latin is None:
                latin = etree.SubElement(rPr, qn("a:latin"))
            latin.set("typeface", "맑은 고딕")
        except Exception:
            pass

        # 제목 하단 구분선 (얇은 직사각형)
        line_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(MARGIN), Inches(TITLE_H + 0.05),
            Inches(SLIDE_W - 2 * MARGIN), Inches(0.03)
        )
        line_bar.fill.solid()
        line_bar.fill.fore_color.rgb = _PRIMARY_ACCENT_RGB
        line_bar.line.fill.background()
        remove_style_element(line_bar._element)
        _remove_shadow(line_bar)

    # 6. 서브그래프 배경 박스 추가 (노드보다 먼저 → 뒤에 위치)
    for sg_idx, sg in enumerate(diagram.subgraphs.values()):
        fill_color = _SUBGRAPH_FILLS[sg_idx % len(_SUBGRAPH_FILLS)]
        _add_subgraph_box(slide, sg, fill_color, sg_idx)

    # 7. 노드 도형 추가 + 인덱스 맵 구성
    shape_map: dict[str, object] = {}
    for node_idx, (nid, node) in enumerate(diagram.nodes.items()):
        # 서브그래프 인덱스를 팔레트 선택에 사용
        if node.subgraph_id and node.subgraph_id in diagram.subgraphs:
            sg_keys = list(diagram.subgraphs.keys())
            palette_idx = sg_keys.index(node.subgraph_id)
        else:
            palette_idx = node_idx

        shape_map[nid] = _add_node_shape(slide, node, palette_idx)

    # 8. 엣지(ELBOW 커넥터 + 화살표) 추가
    for edge in diagram.edges:
        src_shape = shape_map.get(edge.source)
        dst_shape = shape_map.get(edge.target)
        if not src_shape or not dst_shape:
            continue

        _add_connector_elbow(slide, src_shape, dst_shape, label=edge.label)

    # 9. BytesIO로 저장 후 바이트 반환
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
